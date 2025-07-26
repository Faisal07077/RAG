import io
from typing import Dict, Any, List
import PyPDF2
import pdfplumber
from pptx import Presentation
from docx import Document
import pandas as pd
import csv

class DocumentParser:
    """
    Unified document parser for multiple file formats
    """
    
    def __init__(self):
        self.supported_types = ['pdf', 'pptx', 'docx', 'csv', 'text']
    
    def parse(self, file_content: bytes, parser_type: str) -> Dict[str, Any]:
        """Parse document based on type"""
        parser_methods = {
            'pdf': self._parse_pdf,
            'pptx': self._parse_pptx,
            'docx': self._parse_docx,
            'csv': self._parse_csv,
            'text': self._parse_text
        }
        
        if parser_type not in parser_methods:
            raise ValueError(f"Unsupported parser type: {parser_type}")
        
        return parser_methods[parser_type](file_content)
    
    def _parse_pdf(self, file_content: bytes) -> Dict[str, Any]:
        """Parse PDF document"""
        try:
            # Try pdfplumber first (better text extraction)
            with io.BytesIO(file_content) as file_buffer:
                with pdfplumber.open(file_buffer) as pdf:
                    text_parts = []
                    page_count = len(pdf.pages)
                    
                    for page_num, page in enumerate(pdf.pages):
                        page_text = page.extract_text()
                        if page_text:
                            text_parts.append(f"[Page {page_num + 1}]\n{page_text}")
                    
                    return {
                        "text": "\n\n".join(text_parts),
                        "page_count": page_count,
                        "metadata": {"parser": "pdfplumber"}
                    }
        except Exception as e:
            # Fallback to PyPDF2
            try:
                with io.BytesIO(file_content) as file_buffer:
                    pdf_reader = PyPDF2.PdfReader(file_buffer)
                    text_parts = []
                    page_count = len(pdf_reader.pages)
                    
                    for page_num, page in enumerate(pdf_reader.pages):
                        page_text = page.extract_text()
                        if page_text:
                            text_parts.append(f"[Page {page_num + 1}]\n{page_text}")
                    
                    return {
                        "text": "\n\n".join(text_parts),
                        "page_count": page_count,
                        "metadata": {"parser": "PyPDF2"}
                    }
            except Exception as fallback_error:
                raise Exception(f"PDF parsing failed with both pdfplumber and PyPDF2: {str(e)}, {str(fallback_error)}")
    
    def _parse_pptx(self, file_content: bytes) -> Dict[str, Any]:
        """Parse PowerPoint presentation"""
        try:
            with io.BytesIO(file_content) as file_buffer:
                prs = Presentation(file_buffer)
                text_parts = []
                slide_count = len(prs.slides)
                
                for slide_num, slide in enumerate(prs.slides):
                    slide_text_parts = []
                    
                    # Extract text from shapes
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            try:
                                text = shape.text
                                if text and text.strip():
                                    slide_text_parts.append(text.strip())
                            except (AttributeError, Exception):
                                # Skip shapes that don't have accessible text
                                continue
                    
                    if slide_text_parts:
                        slide_text = "\n".join(slide_text_parts)
                        text_parts.append(f"[Slide {slide_num + 1}]\n{slide_text}")
                
                return {
                    "text": "\n\n".join(text_parts),
                    "slide_count": slide_count,
                    "metadata": {"parser": "python-pptx", "slide_count": slide_count}
                }
        except Exception as e:
            raise Exception(f"PPTX parsing failed: {str(e)}")
    
    def _parse_docx(self, file_content: bytes) -> Dict[str, Any]:
        """Parse Word document"""
        try:
            with io.BytesIO(file_content) as file_buffer:
                doc = Document(file_buffer)
                text_parts = []
                
                # Extract paragraphs
                for para in doc.paragraphs:
                    if para.text.strip():
                        text_parts.append(para.text.strip())
                
                # Extract tables
                for table in doc.tables:
                    table_text = []
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            row_text.append(cell.text.strip())
                        table_text.append(" | ".join(row_text))
                    if table_text:
                        text_parts.append("[Table]\n" + "\n".join(table_text))
                
                return {
                    "text": "\n\n".join(text_parts),
                    "paragraph_count": len(doc.paragraphs),
                    "table_count": len(doc.tables),
                    "metadata": {"parser": "python-docx"}
                }
        except Exception as e:
            raise Exception(f"DOCX parsing failed: {str(e)}")
    
    def _parse_csv(self, file_content: bytes) -> Dict[str, Any]:
        """Parse CSV file"""
        try:
            # Try different encodings
            for encoding in ['utf-8', 'latin-1', 'cp1252']:
                try:
                    content_str = file_content.decode(encoding)
                    break
                except UnicodeDecodeError:
                    continue
            else:
                raise Exception("Could not decode CSV file with common encodings")
            
            # Parse CSV
            csv_file = io.StringIO(content_str)
            df = pd.read_csv(csv_file)
            
            # Convert to text representation
            text_parts = []
            
            # Add column headers
            headers = " | ".join(df.columns.tolist())
            text_parts.append(f"[CSV Headers]\n{headers}")
            
            # Add rows (limit to prevent huge text)
            max_rows = min(1000, len(df))
            for i, (idx, row) in enumerate(df.head(max_rows).iterrows()):
                row_text = " | ".join([str(val) for val in row.values])
                text_parts.append(f"Row {i + 1}: {row_text}")
            
            if len(df) > max_rows:
                text_parts.append(f"... and {len(df) - max_rows} more rows")
            
            return {
                "text": "\n".join(text_parts),
                "row_count": len(df),
                "column_count": len(df.columns),
                "metadata": {
                    "parser": "pandas",
                    "columns": df.columns.tolist(),
                    "shape": df.shape
                }
            }
        except Exception as e:
            raise Exception(f"CSV parsing failed: {str(e)}")
    
    def _parse_text(self, file_content: bytes) -> Dict[str, Any]:
        """Parse plain text or markdown file"""
        try:
            # Try different encodings
            for encoding in ['utf-8', 'latin-1', 'cp1252']:
                try:
                    text = file_content.decode(encoding)
                    break
                except UnicodeDecodeError:
                    continue
            else:
                raise Exception("Could not decode text file with common encodings")
            
            return {
                "text": text,
                "line_count": len(text.splitlines()),
                "metadata": {"parser": "text", "encoding": encoding}
            }
        except Exception as e:
            raise Exception(f"Text parsing failed: {str(e)}")
